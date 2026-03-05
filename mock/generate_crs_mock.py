"""Generate a CRS-based mock survey PowerPoint for testing the pipeline.

Creates a realistic 10-story, 19-unit Manhattan condo at 315 East 46th Street
based on reverse-engineered content from an 83-page RAND Capital Reserve Study.
Field notes are "sloppified" to mimic real engineer phone notes taken during
a building walk-through.

- Cover slide
- Section divider slides for 10 systems
- 30 observation slides with real photos extracted from CRS PDF
- 3 intentionally blank slides (after cover, at obs 10, at obs 20)
- 10 distinct note styles for maximum variety
"""

import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = Path(__file__).resolve().parent / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_PATH = OUTPUT_DIR / "test_crs_survey.pptx"
IMAGES_DIR = Path(__file__).resolve().parent / "images"
CROPPED_DIR = IMAGES_DIR / "cropped"

# System colors for placeholder fallback photos
SYSTEM_COLORS = {
    "Roof": (139, 90, 43),
    "Exterior Walls / Facade": (180, 120, 60),
    "Site / Grounds": (34, 139, 34),
    "Common Areas": (160, 140, 180),
    "Interior Finishes": (210, 180, 140),
    "HVAC": (169, 169, 169),
    "Electrical": (255, 165, 0),
    "Plumbing": (70, 130, 180),
    "Fire Protection": (220, 20, 60),
    "Elevators": (192, 192, 192),
}

# Map each observation index (0-29) to a real photo from the CRS PDF.
# Individual photos (588x441) are in images/, cropped composites in images/cropped/.
PHOTO_MAP = [
    # Roof (3 obs)
    "p63_img1_177.jpeg",                    # 0: roof common area / ponding
    "p63_img2_178.jpeg",                    # 1: wood floor boards on roof
    "p64_img4_188.jpeg",                    # 2: parapet metal capping
    # Exterior Walls / Facade (4 obs)
    "p63_img4_180.jpeg",                    # 3: corroding metal detail (facade deterioration)
    "p63_img3_179.jpeg",                    # 4: enclosed mech area / panels
    "p70_img3_212.jpeg",                    # 5: window with leak issues
    "p71_img1_217.jpeg",                    # 6: unit envelope deterioration (sealant related)
    # Site / Grounds (1 obs)
    "p69_img3_206.jpeg",                    # 7: stairway to street level
    # Common Areas (3 obs)
    "cropped/photo_4-3_hallway.jpeg",       # 8: common hallway
    "cropped/photo_4-8_basement_hall.jpeg",  # 9: basement hallway / cellar
    "p82_img3_285.jpeg",                    # 10: trash compactor
    # Interior Finishes (2 obs)
    "p74_img1_238.jpeg",                    # 11: thermostat / energy systems
    "p70_img4_213.jpeg",                    # 12: water damage on wall from window leaks
    # HVAC (6 obs)
    "p72_img3_226.jpeg",                    # 13: BAC cooling tower
    "p73_img4_234.jpeg",                    # 14: boiler
    "p72_img2_225.jpeg",                    # 15: enclosed mech on roof / RTUs
    "p74_img2_239.jpeg",                    # 16: downblast exhaust fans
    "p74_img3_240.jpeg",                    # 17: piping insulation / taping
    "p73_img2_232.jpeg",                    # 18: enclosed mechanical equipment
    # Electrical (4 obs)
    "p75_img3_247.jpeg",                    # 19: electrical panel
    "p76_img3_254.jpeg",                    # 20: rusted conduit / junction
    "p75_img1_245.jpeg",                    # 21: generator
    "p71_img4_220.jpeg",                    # 22: thermostat and intercom
    # Plumbing (3 obs)
    "p77_img3_264.jpeg",                    # 23: domestic water piping
    "p77_img4_265.jpeg",                    # 24: water intake / backflow
    "p69_img2_205.jpeg",                    # 25: basement storage area (cellar level)
    # Fire Protection (2 obs)
    "cropped/photo_9-2_fire_pump.jpeg",     # 26: fire pump
    "p75_img4_248.jpeg",                    # 27: service switch panel (panel-type equipment)
    # Elevators (2 obs)
    "cropped/photo_10-3_machine_room.jpeg",  # 28: elevator machine room
    "cropped/photo_10-2_elev_exterior.jpeg", # 29: elevator cab exterior
]

# Observation data: (system, label_on_photo, sloppy_field_notes)
# 30 observations across 10 systems, 10 distinct note styles
OBSERVATIONS = [
    # ──────────────────────────────────────────────────────────────
    # 1. Roof — 3 observations
    # ──────────────────────────────────────────────────────────────

    # 1-1: Clean SPORT style
    ("Roof", "Roof membrane ponding",
     "S: Main roof, 10th floor level\n"
     "P: Ponding water on modified bitumen membrane, multiple areas\n"
     "O: Membrane is approx 18 yrs old, ponding observed in 3 areas up to 1.5\" deep "
     "after rain 2 days ago. Seams lifting at north side. Base flashings cracked "
     "at parapet. Wood nailers visible where membrane pulled back.\n"
     "R: Commission roof evaluation ($15K) to determine remaining useful life. "
     "Budget full membrane replacement w/ new TPO system incl tapered insulation "
     "for positive drainage.\n"
     "T: Year 1 eval, Year 3-5 replacement\n"
     "Cost: $15,000 evaluation + $364,000 replacement"),

    # 1-2: Quick bullets, no SPORT
    ("Roof", "Combustible wood decking at roof",
     "- wood plank decking under membrane at main roof\n"
     "- not code compliant for new construction but grandfathered\n"
     "- fire retardant treatment unknown, prob original 1960s\n"
     "- if membrane replaced need to address decking\n"
     "- could add coverboard or replace w/ non-combustible\n"
     "- info item for now, coordinate w/ membrane replacement"),

    # 1-3: Partial SPORT
    ("Roof", "Parapet wall condition",
     "S: Parapet walls all elevations, roof level\n"
     "P: Cracking, spalling, deteriorated coping stones\n"
     "O: Through-wall cracks at multiple locations esp NE corner. "
     "Coping stones loose, some missing mortar joints completely open. "
     "Water infiltration likely contributing to top floor apt leaks. "
     "Will need to be addressed as part of FISP cycle 9 facade work.\n"
     "R: Include in FISP close-out / upcoming cycle scope\n"
     "Cost: approx $250K as part of FISP program"),

    # ──────────────────────────────────────────────────────────────
    # 2. Exterior Walls / Facade — 4 observations
    # ──────────────────────────────────────────────────────────────

    # 2-1: Clean SPORT
    ("Exterior Walls / Facade", "Facade condition east elevation",
     "S: East elevation, all floors, brown brick w/ limestone banding\n"
     "P: Spalling brick, cracked lintels, open mortar joints\n"
     "O: FISP cycle 8 close-out completed 2019. Currently in cycle 9 "
     "initial inspection phase. Noted spalling at floors 6-8 east side, "
     "steel lintels corroding at window heads, efflorescence at limestone "
     "bands indicating water migration thru wall assembly.\n"
     "R: Budget for FISP cycle 9 repairs, scaffold full east elevation, "
     "repoint, replace damaged brick, clean and paint lintels\n"
     "T: Year 2-4\n"
     "Cost: $250,000"),

    # 2-2: Stream-of-consciousness with typos
    ("Exterior Walls / Facade", "Plexiglass panels at setback",
     "so theres these plexiglass panels at the 4th floor setback terrace area, "
     "they were installed maybe 10-15 yrs ago as wind screens but theyre all "
     "yellowed and cracked now. a few panels actually missing, just the aluminum "
     "frame left. super ugly from the street. honestly dont know why they used "
     "plexiglass instead of tempered glass in the first place. gonna need full "
     "replacement of all panels w/ new tempered glass or aluminum composite. "
     "the framing is aluminum so thats prob ok just needs cleaning and new "
     "gaskets. estimate around $171K for the whole thing incl scaffolding "
     "access from the setback level. could coordinate w/ facade work to save "
     "on scaffold costs. violaiton potential if panels keep falling off."),

    # 2-3: Minimal (3 lines)
    ("Exterior Walls / Facade", "Water staining at window head",
     "water staining below multiple window heads floors 3-7 east side\n"
     "sealant failed at lintel-to-brick joint, water getting behind facade\n"
     "reseal all window perimeters east elevation ~$15,000"),

    # 2-4: Partial SPORT with heavy abbreviation
    ("Exterior Walls / Facade", "Window sealant deterioration",
     "S: Windows all elevations, typ flrs\n"
     "P: Sealant @ window perimeters deteriorated/missing\n"
     "O: Orig sealant approx 20 yrs old, cracked/pulled away from substrate "
     "at 60%+ of windows checked. Backer rod exposed in several locations. "
     "Air/water infiltration complaints from residents esp upper flrs. "
     "Windows themselves are alum thermal break replacements from ~2005, "
     "frames in decent shape.\n"
     "R: Remove & replace all ext sealant @ window perimeters bldg-wide\n"
     "T: Yr 1-2\n"
     "Cost: ~$50K"),

    # ──────────────────────────────────────────────────────────────
    # 3. Site / Grounds — 1 observation
    # ──────────────────────────────────────────────────────────────

    # 3-1: Info-only note
    ("Site / Grounds", "Sidewalk flags cracked",
     "sidewalk flags along 46th st frontage cracked and heaved in a few spots, "
     "trip hazard near building entrance. DOT likely to flag on next inspection. "
     "building is responsible for sidewalk per NYC admin code. not urgent but "
     "should be on the radar. super will monitor and patch as needed. no cost "
     "assigned for now, minor concrete patching by staff."),

    # ──────────────────────────────────────────────────────────────
    # 4. Common Areas — 3 observations
    # ──────────────────────────────────────────────────────────────

    # 4-1: Clean SPORT
    ("Common Areas", "Corridor flooring worn",
     "S: Common corridors, all floors\n"
     "P: Carpet worn and stained, some areas torn at seams\n"
     "O: Carpet is approx 15 yrs old throughout corridors. Heavy wear "
     "at elevator landings and stairwell entries. Several areas have been "
     "patched with non-matching carpet. Underlayment visible at 3rd and "
     "7th floor elevator lobbies. Walls also showing wear — scuff marks, "
     "dents from move-ins, paint peeling at several corners.\n"
     "R: Full corridor renovation — new flooring (recommend LVT over carpet "
     "for durability), repaint walls, new LED lighting fixtures, update "
     "signage. Phase by floor, 2-3 floors per year.\n"
     "T: Year 1-5 phased\n"
     "Cost: $500,000 total phased"),

    # 4-2: Urgent stream-of-consciousness
    ("Common Areas", "Cellar wall condition",
     "cellar walls are in rough shape — active water seepage at north wall, "
     "can see mineral deposits and staining from ongoing leak. paint is bubbling "
     "and peeling everywhere down here. the sump pump is running constantly which "
     "isnt great (see plumbing notes). theres also some efflorescence on the "
     "CMU walls near the boiler room entry which means water is migrating thru. "
     "honestly need a waterproofing consultant to come look at this, could be "
     "hydrostatic pressure issue or just failed damp-proofing. budget $15K for "
     "evaluation and initial remediation. this needs attention soon before it "
     "gets worse."),

    # 4-3: Bullet list
    ("Common Areas", "Trash compactor condition",
     "- Trash compactor in cellar service area\n"
     "- Unit is Wilkinson model, approx 12 yrs old\n"
     "- Ram operates but slow, hydraulic fluid leak at cylinder\n"
     "- Container rusted at base, drainage clogged\n"
     "- Odor control system not working\n"
     "- Budget $15K for replacement in yr 3-5\n"
     "- Current repairs (hydraulic seal + drain clean): $2,500\n"
     "- Super says it breaks down every few months"),

    # ──────────────────────────────────────────────────────────────
    # 5. Interior Finishes — 2 observations
    # ──────────────────────────────────────────────────────────────

    # 5-1: Mixed format — SPORT then rambling
    ("Interior Finishes", "Building energy label",
     "S: Building-wide, energy systems and envelope\n"
     "P: LL84 benchmarking shows poor energy performance score\n"
     "O: Building energy grade is D per NYC benchmarking. Current EUI around "
     "130 which is high for this bldg type. Main drivers are old boiler plant "
     "(see HVAC section), no building envelope insulation, single-pane windows "
     "at a few remaining original units.\n"
     "anyway the LL97 penalties start hitting in 2024 and this bldg is gonna "
     "be over the threshold. need to start planning decarbonization strategy. "
     "the low hanging fruit is LED lighting conversion in common areas (mostly "
     "done), boiler replacement (see HVAC), and envelope improvements where "
     "feasible.\n"
     "R: Develop LL97 compliance roadmap\n"
     "T: Ongoing\n"
     "Cost: est $150,000 for initial improvements (LED completion + controls)"),

    # 5-2: Dictation style
    ("Interior Finishes", "Apt interior water damage",
     "ok so i looked at unit 6B today per the board's request and theres "
     "definitely water damage at the ceiling near the exterior wall. looks like "
     "its been going on for a while — staining is old and theres some new "
     "bubbling in the paint. the plaster is soft when you push on it. this is "
     "almost certainly from the facade issues noted on the east elevation — "
     "water is getting in at the window head or thru cracks in the brick and "
     "running down behind the plaster. once the facade work is done this should "
     "stop but the interior damage needs repair too. im gonna say about 30K "
     "for interior repairs to affected units on the east side, theres prob "
     "4-5 units with similar issues based on complaint log."),

    # ──────────────────────────────────────────────────────────────
    # 6. HVAC — 6 observations
    # ──────────────────────────────────────────────────────────────

    # 6-1: Clean SPORT
    ("HVAC", "Cooling tower",
     "S: Cooling tower, roof level\n"
     "P: Tower is aging, fill media degraded, drift eliminators damaged\n"
     "O: BAC model, approx 20 yrs old. Fill media sagging, several "
     "sections of drift eliminators cracked/missing. Basin has corrosion "
     "at drain. Motor and fan operational but vibration noted. Water treatment "
     "program in place via Chem-Aqua. Tower services lobby and common area "
     "AC system only (individual units have their own PTACs).\n"
     "R: Replace fill media and drift eliminators now. Budget full tower "
     "replacement in 5-7 yrs.\n"
     "T: Year 1 repairs, Year 5-7 replacement\n"
     "Cost: $15,000 current repairs"),

    # 6-2: Technical, partial SPORT
    ("HVAC", "Boiler plant",
     "S: Boiler room, cellar level\n"
     "P: Dual boiler plant approaching end of useful life\n"
     "O: Two Burnham V906A cast iron boilers, installed approx 2003, natural "
     "gas fired. Boiler #1 has cracked section (repaired but section weeping "
     "slightly). Boiler #2 operating ok but combustion test showed 81% "
     "efficiency which is below what you'd want. Breeching is original "
     "galv steel, rusting at joints. Chimney liner status unknown — recommend "
     "video inspection. Controls are Honeywell but the BMS integration is "
     "minimal, basically just on/off with OAT reset.\n"
     "R: Plan full boiler plant replacement w/ high-eff condensing units, "
     "new breeching, controls upgrade, BMS integration\n"
     "T: Year 3-5\n"
     "Cost: part of $1M mechanical systems budget"),

    # 6-3: Engineer shorthand
    ("HVAC", "Rooftop condensing units",
     "RTU condensing units on roof — 3 units serving common area AC\n"
     "Unit 1: Carrier 50XC, ~15 yr, R-410A, operational, noisy compressor\n"
     "Unit 2: Carrier 50XC, same age, low on charge per super\n"
     "Unit 3: Trane unit, newer ~8 yr, good shape\n"
     "all units on curbs, ductwork thru roof penetrations\n"
     "recommend replacing units 1&2 in next 3-5 yrs\n"
     "coordinate w/ roof membrane replacement\n"
     "part of future mech systems budget"),

    # 6-4: Quick checklist
    ("HVAC", "Exhaust fan enclosure",
     "exhaust fan enclosure on roof:\n"
     "[x] kitchen exhaust fan — operational, belt worn\n"
     "[x] garage exhaust fan — operational, loud bearing\n"
     "[ ] bathroom exhaust riser fans — 2 of 4 not working per residents\n"
     "[x] boiler room exhaust — ok\n"
     "[x] electrical room exhaust — ok\n"
     "need to replace belts on kitchen fan, investigate bathroom exhaust\n"
     "minor maintenance items, super can handle most of it"),

    # 6-5: Messy mixed observations
    ("HVAC", "HVAC piping insulation",
     "piping insulation throughout cellar and mech spaces is a mess. some "
     "sections are original — looks like calcium silicate type, could be "
     "asbestos era material (bldg is 1963). other sections were re-insulated "
     "w/ fiberglass at some point but the jacketing is torn and falling off. "
     "the chilled water lines are sweating where insulation is missing which "
     "is causing corrosion on the pipes and water on the floor.\n"
     "honestly the whole cellar needs an insulation survey. if theres ACM "
     "thats a whole other budget item for abatement.\n"
     "for now just noting as part of overall mech systems budget\n"
     "super says the asbestos survey from 2018 found some in the boiler rm "
     "pipe insulation — need to pull that report"),

    # 6-6: Minimal (3 lines)
    ("HVAC", "Humidifier unit",
     "steam humidifier on AHU serving lobby — not operational\n"
     "DriSteem unit approx 10 yrs, cylinder is shot, control board error\n"
     "replace cylinder and control board $4,500"),

    # ──────────────────────────────────────────────────────────────
    # 7. Electrical — 4 observations
    # ──────────────────────────────────────────────────────────────

    # 7-1: Clean SPORT, code-focused
    ("Electrical", "Electrical panel",
     "S: Main electrical room, cellar\n"
     "P: Main switchgear and distribution panels\n"
     "O: Main service is 1200A 208/120V 3-phase from ConEd vault. GE "
     "Spectra series switchgear, installed approx 2001 during prior "
     "electrical upgrade. Panels in fair condition. Arc flash study was "
     "done in 2019 — labels present but some are faded/illegible. Working "
     "clearances are tight in spots, storage items encroaching on the 36\" "
     "clear zone per NEC 110.26. Need to clean up the room.\n"
     "R: Update arc flash labels, enforce clearance requirements, "
     "thermographic scan recommended\n"
     "T: Year 1\n"
     "Cost: $5,000 for arc flash update and thermo scan"),

    # 7-2: Abbreviated list
    ("Electrical", "Open junction box",
     "found 3 open j-boxes in cellar corridor:\n"
     "1. near boiler rm door — no cover, wires exposed, splice w/o wirenut\n"
     "2. above ceiling at laundry rm — cover missing, abandoned wires\n"
     "3. mech rm wall — cover loose, conduit disconnected\n"
     "code violation NEC 314.25 — all j-boxes must have covers\n"
     "also saw some BX cable w/ damaged jacket near sump pit\n"
     "electrician to correct all — est $7,500\n"
     "priority: safety issue"),

    # 7-3: Detailed technical
    ("Electrical", "Emergency generator",
     "S: Emergency generator, cellar mechanical room\n"
     "P: Diesel generator set for emergency/standby power\n"
     "O: Cummins 100kW diesel genset, installed 2010. Serves emergency "
     "lighting, fire alarm, fire pump, and elevator recall per NYC code. "
     "Auto transfer switch is Asco 300 series. Unit is load bank tested "
     "annually per super — last test report 2024. Fuel oil day tank in "
     "cellar, main storage tank (275 gal) in vault. Generator appeared "
     "clean and well-maintained. Battery charger operational. Block heater "
     "working. Exhaust piping routed to exterior, no visible leaks.\n"
     "R: Continue annual load bank testing and maintenance program. "
     "Budget for major overhaul (injectors, turbo, etc) at 15-yr mark.\n"
     "T: Year 5+ for overhaul\n"
     "Cost: future budget item"),

    # 7-4: Casual conversational
    ("Electrical", "Intercom panel lobby",
     "the intercom system in the lobby is one of those old Aiphone units "
     "from like the early 2000s. it still works but the handsets in a bunch "
     "of units are broken and the directory panel in the vestibule is all "
     "scratched up. the super said they keep replacing handsets but its "
     "getting hard to find parts. honestly for 19 units its not a huge deal "
     "to replace the whole system. the new ones have video and phone integration "
     "which the residents have been asking about anyway. not a capital item "
     "really, just periodic maintenance and eventual upgrade when parts run out."),

    # ──────────────────────────────────────────────────────────────
    # 8. Plumbing — 3 observations
    # ──────────────────────────────────────────────────────────────

    # 8-1: Clean SPORT
    ("Plumbing", "Water booster pump",
     "S: Domestic water booster pump, cellar mechanical room\n"
     "P: Booster pump showing wear, vibration, reduced performance\n"
     "O: Grundfos CRE 15-3 variable speed booster pump, installed 2012. "
     "Pump is vibrating more than normal — bearings likely wearing. Flow "
     "rate adequate but upper floor residents report occasional low pressure "
     "during peak usage. VFD display shows no faults but running at higher "
     "speed than expected for the demand. Pressure tank bladder may also "
     "need replacement — tank is original 2012.\n"
     "R: Rebuild pump (bearings, mechanical seal, impeller inspection). "
     "Replace pressure tank bladder. Evaluate if second pump needed for "
     "redundancy.\n"
     "T: Year 1-2\n"
     "Cost: $25,000"),

    # 8-2: Code-focused abbreviation
    ("Plumbing", "Backflow preventer assembly",
     "BFP in cellar near water entry:\n"
     "- Watts 909 RPZ, 4\", installed ~2015\n"
     "- annual test current per DOB/DEP reqs\n"
     "- relief valve weeping slightly — needs rebuild kit\n"
     "- test cocks operational\n"
     "- strainer upstream needs cleaning\n"
     "- DEP cross-connection report on file\n"
     "est $5K for rebuild and strainer service\n"
     "mandatory compliance item — must maintain"),

    # 8-3: Worried/urgent tone
    ("Plumbing", "Sump pump pit",
     "the sump pump situation in the cellar is concerning. theres one sump "
     "pit with a single pump — no redundancy at all. the pump is running "
     "almost continuously which tells me either groundwater infiltration is "
     "significant or the pump is undersized. the pit is full of sediment, "
     "hasnt been cleaned in who knows how long. if this pump fails during "
     "a heavy rain event the cellar is going to flood. saw water marks on "
     "the walls from previous flooding events — super confirmed it happened "
     "twice last year.\n"
     "strongly recommend installing duplex pump system with alternator and "
     "high water alarm. also clean the pit and install a proper check valve. "
     "this should be year 1 priority.\n"
     "Cost: est $5,000"),

    # ──────────────────────────────────────────────────────────────
    # 9. Fire Protection — 2 observations
    # ──────────────────────────────────────────────────────────────

    # 9-1: Clean SPORT, safety-focused
    ("Fire Protection", "Fire pump",
     "S: Fire pump room, cellar level\n"
     "P: Fire pump operational but components aging\n"
     "O: Aurora 100 GPM electric fire pump, installed 2005. Pump tested "
     "annually per NFPA 25 — last test showed adequate flow and pressure. "
     "However: jockey pump cycles frequently (possible system leak), "
     "controller is showing corrosion at wire terminations, and the test "
     "header drain piping is corroded and leaking slightly. Fire dept "
     "connection (FDC) siamese on facade is functional but caps are missing. "
     "Sprinkler system is wet pipe in cellar and parking only — upper "
     "floors are non-sprinklered (pre-LL26 building).\n"
     "R: Current — repair jockey pump, replace test header piping, new "
     "FDC caps, investigate system leak. Future — evaluate LL26 sprinkler "
     "retrofit requirements and timeline.\n"
     "T: Year 1 current, Year 5-10 sprinkler retrofit\n"
     "Cost: $15,000 current repairs + $125,000 future sprinkler retrofit"),

    # 9-2: Terse, code-violation urgency
    ("Fire Protection", "Fire alarm panel",
     "fire alarm panel — Notifier NFS-320, installed 2008\n"
     "panel showing 3 trouble conditions: 2 supervisory (flow switch zones "
     "3 & 5) and 1 ground fault\n"
     "ground fault needs to be traced and cleared ASAP\n"
     "flow switch issues are prob bad tamper switches — replace\n"
     "also: pull stations in stairwells need testing, havent been tested "
     "this year per super\n"
     "FDNY may cite on next inspection if troubles not cleared\n"
     "minor cost — electrician + fire alarm tech, few thousand $"),

    # ──────────────────────────────────────────────────────────────
    # 10. Elevators — 2 observations
    # ──────────────────────────────────────────────────────────────

    # 10-1: Investigative/questioning
    ("Elevators", "Elevator machine room",
     "looked at the elevator machine room on the roof — its one hydraulic "
     "passenger elevator serving all 10 floors. the machine is a Dover "
     "(now ThyssenKrupp) holeless hydraulic from what looks like a 1990s "
     "modernization. controller is a Motion Control Engineering (MCE) unit. "
     "the machine room is hot — no ventilation fan? need to check code on "
     "that. oil in the tank looks dark, when was it last changed? super "
     "doesnt know. jack packing is seeping a little. the power unit motor "
     "sounds ok but the starter contacts are probably pitted after 30 yrs.\n"
     "the big question is whether this needs a full mod or just component "
     "upgrades. cab is worn too (see next slide). maintenance contract is "
     "with Champion — need to review the contract scope.\n"
     "ongoing maintenance, no specific cost assigned here"),

    # 10-2: Casual planning-oriented
    ("Elevators", "Elevator cab interior",
     "elevator cab interior is pretty tired looking. the pads are always on "
     "because of constant move-ins in a 19-unit building which honestly "
     "makes sense. underneath the pads the walls are scratched stainless, "
     "floor is worn vinyl, ceiling panel is yellowed and one light is out. "
     "the door operator is slow and sometimes catches on peoples bags — "
     "super gets complaints weekly. handrail is loose on one side.\n"
     "the board has been talking about a cab renovation for a while. "
     "recommend doing it as part of a broader elevator modernization that "
     "includes new controller, door operator, cab interior, fixtures, and "
     "ADA-compliant controls. that way they only shut down the elevator "
     "once and do everything together.\n"
     "budget around $250K for full mod including cab\n"
     "T: Year 3-5"),
]


def _make_placeholder_image(label, color, width=800, height=600):
    """Create a colored rectangle with a label as a fallback placeholder photo."""
    img = Image.new("RGB", (width, height), color=color)
    draw = ImageDraw.Draw(img)

    try:
        font = ImageFont.truetype("arial.ttf", 28)
    except (OSError, IOError):
        font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), label, font=font)
    text_w = bbox[2] - bbox[0]
    text_h = bbox[3] - bbox[1]
    x = (width - text_w) // 2
    y = (height - text_h) // 2

    padding = 10
    draw.rectangle(
        [x - padding, y - padding, x + text_w + padding, y + text_h + padding],
        fill=(255, 255, 255, 200)
    )
    draw.text((x, y), label, fill=(0, 0, 0), font=font)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _get_photo(obs_idx, photo_label, system_color):
    """Get the real photo for an observation, falling back to placeholder."""
    if obs_idx < len(PHOTO_MAP):
        photo_rel = PHOTO_MAP[obs_idx]
        # Check cropped/ subfolder first, then main images/ dir
        if photo_rel.startswith("cropped/"):
            photo_path = CROPPED_DIR / photo_rel.split("/", 1)[1]
        else:
            photo_path = IMAGES_DIR / photo_rel
        if photo_path.exists():
            buf = io.BytesIO(photo_path.read_bytes())
            return buf
        else:
            print(f"  WARNING: Photo not found: {photo_path}")
    return _make_placeholder_image(photo_label, system_color)


def generate():
    """Generate the CRS-based mock survey PowerPoint."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # --- Cover slide ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Capital Reserve Study \u2014 Physical Condition Survey"
    p.font.size = Pt(36)
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "315 East 46th Street"
    p2.font.size = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "New York, NY 10017"
    p3.font.size = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "RAND Engineering & Architecture, DPC"
    p4.font.size = Pt(18)

    p5 = tf.add_paragraph()
    p5.text = "March 2026"
    p5.font.size = Pt(18)

    # --- Blank slide #1 (intentional, after cover) ---
    prs.slides.add_slide(blank_layout)

    # Track current system to insert dividers
    current_system = None
    obs_idx = 0
    real_photos = 0
    fallback_photos = 0

    for i, (system, photo_label, notes) in enumerate(OBSERVATIONS):
        # Insert section divider when system changes
        if system != current_system:
            current_system = system
            slide = prs.slides.add_slide(blank_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
            tf = txBox.text_frame
            sys_num = list(SYSTEM_COLORS.keys()).index(system) + 1
            p = tf.paragraphs[0]
            p.text = f"{sys_num}. {system}"
            p.font.size = Pt(32)
            p.font.bold = True

        # Insert blank slides at observation 10 and 20
        if obs_idx == 10:
            prs.slides.add_slide(blank_layout)
        if obs_idx == 20:
            prs.slides.add_slide(blank_layout)

        # --- Observation slide ---
        slide = prs.slides.add_slide(blank_layout)
        color = SYSTEM_COLORS.get(system, (128, 128, 128))

        # Add real photo (or fallback placeholder)
        img_buf = _get_photo(obs_idx, photo_label, color)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(0.5), Inches(6), Inches(4.5))

        # Track photo source
        if obs_idx < len(PHOTO_MAP):
            photo_rel = PHOTO_MAP[obs_idx]
            check = CROPPED_DIR / photo_rel.split("/", 1)[1] if photo_rel.startswith("cropped/") else IMAGES_DIR / photo_rel
            if check.exists():
                real_photos += 1
            else:
                fallback_photos += 1
        else:
            fallback_photos += 1

        # Add field notes
        txBox = slide.shapes.add_textbox(Inches(7), Inches(0.5), Inches(5.5), Inches(6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = notes
        p.font.size = Pt(12)

        obs_idx += 1

    prs.save(str(OUTPUT_PATH))
    n_systems = len(set(s for s, _, _ in OBSERVATIONS))
    n_obs = len(OBSERVATIONS)
    n_dividers = n_systems
    n_blanks = 3  # after cover, at obs 10, at obs 20
    total = 1 + n_blanks + n_dividers + n_obs
    print(f"CRS mock survey created: {OUTPUT_PATH}")
    print(f"  - 1 cover slide")
    print(f"  - {n_dividers} section dividers")
    print(f"  - {n_obs} observation slides ({real_photos} real photos, {fallback_photos} placeholders)")
    print(f"  - {n_blanks} blank slides")
    print(f"  - {total} total slides")


if __name__ == "__main__":
    generate()
