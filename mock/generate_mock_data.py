"""Generate a mock survey PowerPoint for testing the pipeline.

Creates a fictional 12-story co-op at 250 West 72nd Street with:
- Cover slide
- Section divider slides for ~10 systems
- 20 observation slides with placeholder photos + SPORT notes
- 2-3 intentionally blank slides
- SPORT notes of varying quality (clean, abbreviated, messy)
"""

import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = Path(__file__).resolve().parent / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
OUTPUT_PATH = OUTPUT_DIR / "test_survey.pptx"

# System colors for placeholder photos
SYSTEM_COLORS = {
    "Roof": (139, 90, 43),
    "Exterior Walls / Facade": (180, 120, 60),
    "Windows": (100, 149, 237),
    "Plumbing": (70, 130, 180),
    "HVAC": (169, 169, 169),
    "Electrical": (255, 165, 0),
    "Elevators": (192, 192, 192),
    "Fire Protection": (220, 20, 60),
    "Interior Finishes": (210, 180, 140),
    "Site / Grounds": (34, 139, 34),
}

# Observation data: (system, label_on_photo, sport_notes)
OBSERVATIONS = [
    # Roof — 3 observations
    ("Roof", "Roof membrane blistering",
     "S: Main roof, level 12\nP: EPDM membrane blistering and ridging, multiple patches\nO: Membrane approx 20 yrs old, blisters up to 12\" dia, base flashings pulling away at parapet\nR: Full membrane replacement with new TPO, replace base flashings and sheet metal\nT: Year 1-2\nCost: $150,000-$200,000"),
    ("Roof", "Roof drain and scupper",
     "S: Main roof drains\nP: Clogged roof drains, ponding water around drains\nO: 3 of 6 roof drains partially blocked with debris, standing water up to 2\" deep\nR: Clean all drains, install new drain guards, re-slope as needed\nT: Year 1-2"),
    ("Roof", "Bulkhead door deteriorated",
     "bulkhead door rusted out at bottom, doesn't close properly, water getting in\nreplace door\nyear 1-2"),

    # Exterior Walls — 3 observations
    ("Exterior Walls / Facade", "Brick spalling east facade",
     "S: East facade, floors 8-10\nP: Brick spalling and mortar deterioration, steel lintel corrosion\nO: Approx 200 SF affected, lintels rusting and jacking bricks, loose bricks noted\nR: Scaffold east facade, replace brick and mortar, clean and paint lintels\nT: Year 1-2"),
    ("Exterior Walls / Facade", "Parapet cracking",
     "S: North parapet\nP: Through-wall cracks in parapet, cap flashing loose\nO: Multiple cracks noted, some with efflorescence indicating water penetration\nR: Rebuild parapet, install new cap flashing and through-wall flashing\nT: Year 1-2"),
    ("Exterior Walls / Facade", "Sealant at expansion joint",
     "expansion joint sealant failing on west side btwn floors 3-5, old caulk cracked and peeling\nneed to remove and replace all sealant\nmaybe year 3-5"),

    # Windows — 2 observations
    ("Windows", "Window frames corroded",
     "S: Steel casement windows, typical floors\nP: Frames corroded, glazing compound failed, air infiltration\nO: Original 1962 steel casements, single glazed, frames showing surface corrosion, putty glazing cracked and missing in spots\nR: Full window replacement with thermally broken aluminum, insulated glass\nT: Year 3-5"),
    ("Windows", "Lobby entrance door",
     "lobby entrance doors — glass cracked in one panel, hardware worn, closer not working right\nreplace both entrance doors with new aluminum storefront\nyear 1-2"),

    # Plumbing — 2 observations
    ("Plumbing", "Galvanized water risers",
     "S: Domestic water risers, basement mech room\nP: Galv steel risers heavy corrosion, pinhole leaks, clamp repairs\nO: Original 1962 galvanized risers, multiple clamp repairs, reduced flow on upper floors\nR: Replace all galvanized risers with Type L copper\nT: Year 1-3\nCost: $300,000-$400,000"),
    ("Plumbing", "Water heater",
     "S: Domestic hot water heater\nP: Unit past useful life, inefficient, showing corrosion at base\nO: Gas-fired storage type, 500 gal, approx 25 yrs old, flue showing corrosion\nR: Replace with new high-efficiency condensing unit\nT: Year 1-2"),

    # HVAC — 2 observations
    ("HVAC", "Boiler plant",
     "S: Heating boiler, basement boiler room\nP: Cast iron sectional boiler, obsolete controls, inefficient\nO: Weil-McLain Model 88, approx 1985, flame safeguard controls outdated, combustion efficiency measured at 78%\nR: Replace with new high-efficiency condensing boiler, new controls, new breeching\nT: Year 3-5"),
    ("HVAC", "Unit ventilator typical apt",
     "through-wall AC sleeves rusting in multiple units, condensate leaking onto facade\nrecommend new PTAC units throughout\nyear 3-5 probably"),

    # Electrical — 2 observations
    ("Electrical", "Main switchgear",
     "S: Main electrical switchgear, basement electrical room\nP: Original equipment, obsolete, parts unavailable\nO: GE AKD-5 switchgear circa 1962, 1200A service, no arc flash labels, clearances tight\nR: Replace main switchgear, new CT cabinet, coordinate with utility\nT: Year 1-3"),
    ("Electrical", "Emergency lighting",
     "S: Emergency/egress lighting\nP: Insufficient emergency lighting in corridors and stairs\nO: Multiple units non-functional, battery packs dead, does not meet current code\nR: Install new LED emergency lighting throughout common areas per NYC BC\nT: Year 1-2"),

    # Elevators — 1 observation
    ("Elevators", "Elevator machine room",
     "S: Passenger elevators 1 and 2\nP: Frequent breakdowns, long wait times, obsolete controllers\nO: Original 1965 Otis geared traction, relay logic controls, slow door operators\nR: Full modernization — new controllers, door operators, cab interiors, fixtures\nT: Year 3-5\nCost: $500,000-$650,000"),

    # Fire Protection — 1 observation
    ("Fire Protection", "Standpipe system",
     "S: Standpipe and sprinkler system\nP: Standpipe siamese connections corroded, FD connections non-functional\nO: 2 of 3 siamese connections seized, hose valves in stairwells need new caps and handles\nR: Replace all siamese connections, rebuild hose valves, hydro test system\nT: Year 1-2"),

    # Interior Finishes — 2 observations
    ("Interior Finishes", "Lobby floor and walls",
     "lobby floor cracked tiles, walls need paint, ceiling water stain near entrance, mailboxes dented\nrecommend retile lobby floor, paint, replace ceiling tiles, new mailboxes\nyear 1-2"),
    ("Interior Finishes", "Stairwell condition",
     "S: Stairwell A, all floors\nP: Worn treads, loose handrails, peeling paint\nO: Rubber treads worn smooth on floors 1-5, metal handrails loose at multiple brackets, paint peeling throughout\nR: Replace stair treads, secure handrails, prep and repaint all stairwell surfaces\nT: Year 1-2"),

    # Site / Grounds — 2 observations
    ("Site / Grounds", "Sidewalk and areaway",
     "S: Public sidewalk and building areaway\nP: Sidewalk flags heaved and cracked, trip hazards, areaway grating rusted\nO: Multiple trip hazards exceeding DOT threshold, cellar areaway grating corroded\nR: Replace sidewalk flags per DOT, new areaway gratings, waterproof areaway walls\nT: Year 1-2"),
    ("Site / Grounds", "Courtyard pavement",
     "courtyard pavers settling and uneven, drainage backing up in rain, planters falling apart\nregrade and repave courtyard, fix drainage, rebuild planters\nyear 3-5"),
]


def _make_placeholder_image(label, color, width=800, height=600):
    """Create a colored rectangle with a label as a placeholder photo."""
    img = Image.new("RGB", (width, height), color=color)
    draw = ImageDraw.Draw(img)

    # Draw label text
    try:
        font = ImageFont.truetype("arial.ttf", 28)
    except (OSError, IOError):
        font = ImageFont.load_default()

    # Center text
    bbox = draw.textbbox((0, 0), label, font=font)
    text_w = bbox[2] - bbox[0]
    text_h = bbox[3] - bbox[1]
    x = (width - text_w) // 2
    y = (height - text_h) // 2

    # Draw white background rect for readability
    padding = 10
    draw.rectangle(
        [x - padding, y - padding, x + text_w + padding, y + text_h + padding],
        fill=(255, 255, 255, 200)
    )
    draw.text((x, y), label, fill=(0, 0, 0), font=font)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def generate():
    """Generate the mock survey PowerPoint."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # --- Cover slide ---
    slide = prs.slides.add_slide(blank_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Physical Condition Survey"
    p.font.size = Pt(36)
    p.font.bold = True

    p2 = tf.add_paragraph()
    p2.text = "250 West 72nd Street"
    p2.font.size = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "New York, NY 10023"
    p3.font.size = Pt(20)

    p4 = tf.add_paragraph()
    p4.text = "RAND Engineering & Architecture, DPC"
    p4.font.size = Pt(18)

    p5 = tf.add_paragraph()
    p5.text = "March 2026"
    p5.font.size = Pt(18)

    # --- Blank slide #1 (intentional) ---
    prs.slides.add_slide(blank_layout)

    # Track current system to insert dividers
    current_system = None
    obs_idx = 0

    for i, (system, photo_label, notes) in enumerate(OBSERVATIONS):
        # Insert section divider when system changes
        if system != current_system:
            current_system = system
            slide = prs.slides.add_slide(blank_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
            tf = txBox.text_frame
            # Find system number
            sys_num = i + 1
            for cat_num, (sys_name, _) in enumerate(SYSTEM_COLORS.items(), 1):
                if sys_name == system:
                    sys_num = cat_num
                    break
            p = tf.paragraphs[0]
            p.text = f"{sys_num}. {system}"
            p.font.size = Pt(32)
            p.font.bold = True

        # Insert blank slide after observation 7 and 14
        if obs_idx == 7:
            prs.slides.add_slide(blank_layout)
        if obs_idx == 14:
            prs.slides.add_slide(blank_layout)

        # --- Observation slide ---
        slide = prs.slides.add_slide(blank_layout)
        color = SYSTEM_COLORS.get(system, (128, 128, 128))

        # Add placeholder photo
        img_buf = _make_placeholder_image(photo_label, color)
        slide.shapes.add_picture(img_buf, Inches(0.5), Inches(0.5), Inches(6), Inches(4.5))

        # Add SPORT notes
        txBox = slide.shapes.add_textbox(Inches(7), Inches(0.5), Inches(5.5), Inches(6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = notes
        p.font.size = Pt(12)

        obs_idx += 1

    prs.save(str(OUTPUT_PATH))
    print(f"Mock survey created: {OUTPUT_PATH}")
    print(f"  - 1 cover slide")
    print(f"  - {len(set(s for s, _, _ in OBSERVATIONS))} section dividers")
    print(f"  - {len(OBSERVATIONS)} observation slides")
    print(f"  - 3 blank slides")
    total = 1 + len(set(s for s, _, _ in OBSERVATIONS)) + len(OBSERVATIONS) + 3
    print(f"  - {total} total slides")


if __name__ == "__main__":
    generate()
