"""Tests for the PowerPoint parser."""

import io
import os
import sys
import tempfile

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

# Ensure src is importable
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from src.parser import parse_pptx, _match_section, _get_slide_text


def _make_test_image():
    """Create a minimal test PNG image as bytes."""
    img = Image.new("RGB", (100, 100), color=(255, 0, 0))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf


def _create_test_pptx(slides_config):
    """Create a test .pptx from a list of slide configs.

    Each config is a dict with optional keys:
      - text: str or None
      - image: bool (whether to add a placeholder image)
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    for config in slides_config:
        slide = prs.slides.add_slide(blank_layout)

        if config.get("text"):
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
            tf = txBox.text_frame
            # Split on \n to create separate paragraphs (python-pptx converts
            # \n within a paragraph to \x0b, so we use add_paragraph instead)
            lines = config["text"].split("\n")
            tf.paragraphs[0].text = lines[0]
            for line in lines[1:]:
                p = tf.add_paragraph()
                p.text = line

        if config.get("image"):
            img_buf = _make_test_image()
            slide.shapes.add_picture(img_buf, Inches(1), Inches(3), Inches(2), Inches(2))

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return tmp.name


class TestMatchSection:
    def test_matches_known_system(self):
        name, is_div = _match_section("1. Roof")
        assert name == "Roof"
        assert is_div is True

    def test_matches_without_number(self):
        name, is_div = _match_section("Roof")
        assert name == "Roof"
        assert is_div is True

    def test_matches_alias(self):
        name, is_div = _match_section("Facade")
        assert name == "Exterior Walls / Facade"
        assert is_div is True

    def test_no_match_long_text(self):
        name, is_div = _match_section(
            "The roof membrane was observed to be in poor condition with blistering."
        )
        assert is_div is False

    def test_empty_text(self):
        name, is_div = _match_section("")
        assert name is None
        assert is_div is False

    def test_numbered_unknown_system(self):
        name, is_div = _match_section("15. Custom System")
        assert name == "Custom System"
        assert is_div is True


class TestParsePptx:
    def test_blank_slides_skipped(self):
        path = _create_test_pptx([
            {"text": "Cover Slide"},  # cover
            {},  # blank — should be skipped
            {"text": "1. Roof"},  # divider
            {"text": "S: Main roof\nP: Blistering", "image": True},  # observation
            {},  # blank — should be skipped
            {"text": "S: Second obs", "image": True},  # observation
        ])
        try:
            report = parse_pptx(path)
            assert len(report.observations) == 2
        finally:
            os.unlink(path)

    def test_cover_slide_extracted(self):
        path = _create_test_pptx([
            {"text": "250 West 72nd Street\nNew York, NY"},
            {"text": "1. Roof"},
            {"text": "Observation text", "image": True},
        ])
        try:
            report = parse_pptx(path)
            assert report.building_name == "250 West 72nd Street"
        finally:
            os.unlink(path)

    def test_section_divider_sets_section(self):
        path = _create_test_pptx([
            {"text": "Building Name"},  # cover
            {"text": "1. Roof"},  # divider
            {"text": "Roof obs text", "image": True},  # observation
            {"text": "5. Plumbing"},  # divider
            {"text": "Plumbing obs text", "image": True},  # observation
        ])
        try:
            report = parse_pptx(path)
            assert len(report.observations) == 2
            assert report.observations[0].section_name == "Roof"
            assert report.observations[1].section_name == "Plumbing"
        finally:
            os.unlink(path)

    def test_image_only_slide_is_observation(self):
        path = _create_test_pptx([
            {"text": "Building"},  # cover
            {"text": "1. Roof"},  # divider
            {"image": True},  # image only — still an observation
        ])
        try:
            report = parse_pptx(path)
            assert len(report.observations) == 1
            assert report.observations[0].image_bytes is not None
        finally:
            os.unlink(path)

    def test_text_only_slide_is_observation(self):
        path = _create_test_pptx([
            {"text": "Building"},  # cover
            {"text": "1. Roof"},  # divider
            {"text": "S: Main roof area\nP: Membrane is deteriorating badly\nO: Multiple patches\nR: Replace membrane\nT: Year 1-2"},
        ])
        try:
            report = parse_pptx(path)
            assert len(report.observations) == 1
            assert report.observations[0].raw_text is not None
        finally:
            os.unlink(path)

    def test_obs_numbers_sequential_within_section(self):
        path = _create_test_pptx([
            {"text": "Building"},  # cover
            {"text": "1. Roof"},  # divider
            {"text": "First obs", "image": True},
            {"text": "Second obs", "image": True},
            {"text": "Third obs", "image": True},
        ])
        try:
            report = parse_pptx(path)
            assert len(report.observations) == 3
            assert report.observations[0].obs_number == "1-1"
            assert report.observations[1].obs_number == "1-2"
            assert report.observations[2].obs_number == "1-3"
        finally:
            os.unlink(path)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
