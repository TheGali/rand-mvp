"""PowerPoint parser for RAND survey presentations.

Extracts observations from .pptx files. Silently skips blank slides.
Detects section dividers to group observations by building system.
"""

import io
import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.models import Observation, Report

# Load system categories for section matching
_CATEGORIES_PATH = Path(__file__).resolve().parent.parent / "style_guide" / "system_categories.json"

try:
    with open(_CATEGORIES_PATH, encoding="utf-8") as f:
        SYSTEM_CATEGORIES = json.load(f)
except FileNotFoundError:
    SYSTEM_CATEGORIES = []


def _get_slide_text(slide):
    """Extract all text from a slide, concatenated."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return "\n".join(texts)


def _get_slide_image(slide):
    """Extract the first image from a slide as bytes, or None."""
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape.image.blob
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.has_text_frame is False:
            try:
                return shape.image.blob
            except (AttributeError, ValueError):
                pass
    return None


def _match_section(text):
    """Try to match slide text to a known system category.

    Returns (section_name, is_divider) tuple.
    is_divider=True means this slide is a section header, not an observation.
    """
    clean = text.strip()
    if not clean:
        return None, False

    # Check for numbered pattern like "1. Roof" or "1 - Roof"
    # Limit to 1-2 digit numbers to avoid matching street addresses (e.g. "250 West 72nd")
    numbered = re.match(r"^\d{1,2}[\.\)\-\s]+(.+)$", clean)
    candidate = numbered.group(1).strip() if numbered else clean

    # Short text (likely a divider) — match against categories
    lines = clean.split("\n")
    if len(lines) <= 2 and len(clean) < 80:
        for cat in SYSTEM_CATEGORIES:
            if candidate.lower() == cat["name"].lower():
                return cat["name"], True
            for alias in cat.get("aliases", []):
                if candidate.lower() == alias.lower():
                    return cat["name"], True
        # If it's short and numbered but didn't match a category, still treat as divider
        if numbered:
            return candidate, True

    return None, False


def parse_pptx(file_path):
    """Parse a PowerPoint file into a Report with Observations.

    Args:
        file_path: Path to the .pptx file

    Returns:
        Report with populated observations list
    """
    prs = Presentation(file_path)
    report = Report()
    current_section = "General"
    section_obs_count = {}
    section_number = 0

    for slide_idx, slide in enumerate(prs.slides):
        text = _get_slide_text(slide)
        image_bytes = _get_slide_image(slide)

        # Skip truly empty slides
        if not text and not image_bytes:
            continue

        # Check for cover slide (first slide with substantial text but no image)
        if slide_idx == 0 and not image_bytes and text:
            # Try to extract building info from cover
            lines = text.split("\n")
            if lines:
                report.building_name = lines[0]
            if len(lines) > 1:
                report.address = lines[1]
            continue

        # Check if this is a section divider
        section_name, is_divider = _match_section(text)
        if is_divider and section_name:
            current_section = section_name
            section_number += 1
            if current_section not in section_obs_count:
                section_obs_count[current_section] = 0
            continue

        # Skip slides that look like placeholders / non-observation content
        # (very short text, no image, and not matching a section)
        if not image_bytes and text and len(text) < 15 and not section_name:
            continue

        # This is an observation slide
        if current_section not in section_obs_count:
            section_obs_count[current_section] = 0
        section_obs_count[current_section] += 1

        # Find the section number for this section
        sec_num = 0
        for cat in SYSTEM_CATEGORIES:
            if cat["name"] == current_section:
                sec_num = cat["number"]
                break
        if sec_num == 0:
            sec_num = section_number if section_number > 0 else 1

        obs_number = f"{sec_num}-{section_obs_count[current_section]}"

        obs = Observation(
            slide_index=slide_idx,
            obs_number=obs_number,
            section_name=current_section,
            raw_text=text,
            image_bytes=image_bytes,
        )
        report.observations.append(obs)

    return report
