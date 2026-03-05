"""Write AI-generated results back into a PowerPoint file.

Adds captions and prose to observation slides. Saves as a new file.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def write_results_to_pptx(report, input_path, output_path):
    """Write processed results back into a copy of the input PowerPoint.

    Args:
        report: Report with populated processed_observations
        input_path: Path to original .pptx
        output_path: Path for the output .pptx
    """
    prs = Presentation(input_path)

    # Build lookup: slide_index -> processed observation
    obs_by_slide = {po.slide_index: po for po in report.processed_observations}

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx not in obs_by_slide:
            continue

        po = obs_by_slide[slide_idx]

        # Add caption text box at the bottom-left
        caption_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.2), Inches(6), Inches(0.5)
        )
        tf = caption_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = po.caption
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)

        # Add prose text box below the notes area
        prose_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5.8), Inches(12), Inches(1.5)
        )
        tf = prose_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = po.prose
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 0, 0)

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"  Written: {output_path}")
