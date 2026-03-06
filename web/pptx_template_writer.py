"""Generate a blank survey template PPTX that satisfies all parser invariants.

Engineers fill in this template with photos and SPORT notes, then upload it.
The structure (cover, dividers, observation placeholders) is pre-built so the
parser never hits structural pitfalls.
"""

import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from web.pptx_report_writer import (
    SLIDE_W, SLIDE_H, NAVY, WHITE, DARK, GRAY, LIGHT_BG,
    _add_text_box, _set_slide_bg,
)

CATEGORIES_PATH = Path(__file__).resolve().parent.parent / "style_guide" / "system_categories.json"

SPORT_PLACEHOLDER = (
    "S: [System]\n"
    "P: [Position \u2014 where in the building]\n"
    "O: [Observation \u2014 what you see, materials, quantities, age]\n"
    "R: [Recommendation \u2014 what to do about it]\n"
    "T: [Timeframe \u2014 e.g. 1 to 3 years]"
)

COVER_NOTES = (
    "INSTRUCTIONS\n"
    "1. Replace [Building Name] with the actual building name.\n"
    "2. Replace [Street Address, City, State ZIP] with the full address.\n"
    "3. Replace [Month Year] with the survey date (e.g. March 2026).\n"
    "4. Do NOT add images to this cover slide."
)

DIVIDER_NOTES = (
    "SECTION DIVIDER \u2014 do not modify the text on this slide.\n"
    "Add observation slides AFTER this divider, one photo per slide."
)

OBS_NOTES = (
    "OBSERVATION SLIDE\n"
    "1. Replace the gray box with a single field photo.\n"
    "2. Fill in the SPORT fields on the right.\n"
    "3. One observation per slide. Duplicate this slide for more observations.\n"
    "4. Keep all observation slides between their section divider and the next."
)


def _load_categories():
    with open(CATEGORIES_PATH, "r", encoding="utf-8") as f:
        return json.load(f)


def _add_speaker_notes(slide, text):
    from pptx.oxml.ns import qn
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def _add_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, NAVY)

    _add_text_box(
        slide, Inches(1), Inches(1.8), Inches(11), Inches(1),
        "[Building Name]",
        font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, Inches(1), Inches(3.0), Inches(11), Inches(0.6),
        "[Street Address, City, State ZIP]",
        font_size=18, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
        "RAND Engineering & Architecture, DPC",
        font_size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER,
    )
    _add_text_box(
        slide, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
        "[Month Year]",
        font_size=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER,
    )

    _add_speaker_notes(slide, COVER_NOTES)
    return slide


def _add_divider(prs, number, name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, LIGHT_BG)

    _add_text_box(
        slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
        f"{number}. {name}",
        font_size=32, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
    )

    _add_speaker_notes(slide, DIVIDER_NOTES)
    return slide


def _add_placeholder_observation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Gray rectangle as photo placeholder
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.5), Inches(5.5), Inches(4.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xD1, 0xD5, 0xDB)
    shape.line.fill.background()  # no border

    # "Insert photo here" label
    _add_text_box(
        slide, Inches(1.5), Inches(2.4), Inches(3.5), Inches(0.5),
        "[ Insert photo here ]",
        font_size=14, color=GRAY, align=PP_ALIGN.CENTER,
    )

    # SPORT text box on the right
    txBox = slide.shapes.add_textbox(
        Inches(6.5), Inches(0.5), Inches(6.3), Inches(5.0),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = SPORT_PLACEHOLDER
    for para in tf.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = DARK
        para.space_after = Pt(6)

    _add_speaker_notes(slide, OBS_NOTES)
    return slide


def generate_template_pptx() -> bytes:
    """Build a blank survey template and return the file bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    categories = _load_categories()

    # Slide 0: cover
    _add_cover(prs)

    # Slides 1..24: 12 dividers, each followed by 1 placeholder observation
    for cat in categories:
        _add_divider(prs, cat["number"], cat["name"])
        _add_placeholder_observation(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
