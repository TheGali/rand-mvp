"""Polished PPTX report generator — the primary deliverable.

Structure:
  1. Cover slide (building name, address, date)
  2. Command center table(s) with hyperlinks to observation slides
  3. Observation slides (one per obs: photo + caption + prose + recommendation)
  4. Funding summary table
"""

import io
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from lxml import etree


# RAND brand colors
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
BLUE = RGBColor(0x25, 0x63, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1F, 0x2A, 0x37)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)

# Slide dimensions (widescreen 13.33 x 7.5 inches)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Condition color map
CONDITION_COLORS = {
    "critical": RGBColor(0x7F, 0x1D, 0x1D),
    "poor":     RGBColor(0x99, 0x1B, 0x1B),
    "fair":     RGBColor(0x92, 0x40, 0x0E),
    "good":     RGBColor(0x06, 0x5F, 0x46),
    "satisfactory": RGBColor(0x1E, 0x40, 0xAF),
}

CONDITION_BG = {
    "critical": RGBColor(0xFE, 0xE2, 0xE2),
    "poor":     RGBColor(0xFE, 0xCA, 0xCA),
    "fair":     RGBColor(0xFE, 0xF3, 0xC7),
    "good":     RGBColor(0xD1, 0xFA, 0xE5),
    "satisfactory": RGBColor(0xDB, 0xEA, 0xFE),
}

ROWS_PER_CC_SLIDE = 15


def _add_text_box(slide, left, top, width, height, text, font_size=11,
                  bold=False, italic=False, color=DARK, align=PP_ALIGN.LEFT):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def _set_slide_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_cover_slide(prs, job):
    """Slide 1: Cover with building name, address, date."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    _set_slide_bg(slide, NAVY)

    # Building name
    _add_text_box(
        slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
        job.get("building_name", "Building Survey"),
        font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )

    # Address
    _add_text_box(
        slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
        job.get("address", ""),
        font_size=18, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER,
    )

    # Report type
    _add_text_box(
        slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
        "Physical Condition Survey",
        font_size=16, italic=True, color=RGBColor(0x99, 0xAA, 0xBB), align=PP_ALIGN.CENTER,
    )

    # RAND branding
    _add_text_box(
        slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
        "RAND Engineering & Architecture, DPC",
        font_size=12, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER,
    )

    return slide


def _add_command_center_slide(prs, observations, slide_start_idx, obs_slides_map):
    """Add a command center table slide with hyperlinks to observation slides."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
        "Command Center", font_size=20, bold=True, color=NAVY,
    )

    # Table
    cols = 7
    rows = len(observations) + 1  # header + data
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(12.3)
    height = Inches(0.35) * rows

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column widths
    col_widths = [Inches(0.7), Inches(2.0), Inches(1.5), Inches(1.2), Inches(1.2), Inches(4.5), Inches(1.2)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    headers = ["Obs #", "System", "Component", "Condition", "Priority", "Summary", "Cost High"]
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data rows
    for i, obs in enumerate(observations):
        row_idx = i + 1
        obs_num = obs.get("obs_number", "")
        values = [
            obs_num,
            obs.get("system", ""),
            obs.get("component", ""),
            obs.get("condition", ""),
            obs.get("priority", ""),
            (obs.get("prose", "")[:80] + "...") if len(obs.get("prose", "")) > 80 else obs.get("prose", ""),
            f"${obs.get('cost_high', 0):,.0f}" if obs.get("cost_high", 0) > 0 else "—",
        ]

        for j, val in enumerate(values):
            cell = table.cell(row_idx, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = DARK

            # Alternating row color
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

        # Condition cell color
        cond = obs.get("condition", "").lower()
        if cond in CONDITION_COLORS:
            cond_cell = table.cell(row_idx, 3)
            for paragraph in cond_cell.text_frame.paragraphs:
                paragraph.font.color.rgb = CONDITION_COLORS[cond]
                paragraph.font.bold = True

        # Add hyperlink from Obs # cell to observation slide
        if obs_num in obs_slides_map:
            target_slide_idx = obs_slides_map[obs_num]
            _add_inter_slide_hyperlink(table.cell(row_idx, 0), prs, target_slide_idx)

    return slide


def _add_inter_slide_hyperlink(cell, prs, target_slide_idx):
    """Add a hyperlink from a table cell to another slide in the presentation.

    Uses python-pptx XML layer for inter-slide action settings.
    """
    try:
        # Get the slide's rId
        target_slide = prs.slides[target_slide_idx]
        # Access the paragraph run
        paragraph = cell.text_frame.paragraphs[0]
        if not paragraph.runs:
            return

        run = paragraph.runs[0]
        run.font.color.rgb = BLUE
        run.font.underline = True

        # Get relationship to target slide
        slide = cell._tc.getparent().getparent().getparent().getparent()
        # Find the slide part
        for rel in prs.part.rels.values():
            if hasattr(rel, 'target_part') and rel.target_part == target_slide.part:
                rId = rel.rId
                break
        else:
            # Need to add a relationship
            rId = prs.part.relate_to(target_slide.part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')

        # Add hlinkClick to the run
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                 'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}
        rPr = run._r.get_or_add_rPr()
        hlinkClick = etree.SubElement(rPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}hlinkClick')
        hlinkClick.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id', rId)
        hlinkClick.set('action', 'ppaction://hlinksldjump')
    except Exception:
        pass  # Hyperlink is best-effort; don't break report generation


def _add_observation_slide(prs, obs, job_dir):
    """Add an observation slide: photo + caption + prose + recommendation."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    obs_number = obs.get("obs_number", "")
    system = obs.get("system", "")

    # Header bar
    _add_text_box(
        slide, Inches(0.4), Inches(0.2), Inches(12), Inches(0.4),
        f"Observation {obs_number} — {system}",
        font_size=16, bold=True, color=NAVY,
    )

    # Try to add photo
    image_path = Path(job_dir) / "images" / f"obs_{obs_number}.png"
    photo_bottom = Inches(0.7)
    if image_path.exists():
        try:
            slide.shapes.add_picture(
                str(image_path), Inches(0.5), Inches(0.7), width=Inches(5.5)
            )
            photo_bottom = Inches(5.0)  # approximate; actual height varies
        except Exception:
            photo_bottom = Inches(0.7)

    # Caption (under photo or at top-right)
    _add_text_box(
        slide, Inches(0.5), photo_bottom, Inches(5.5), Inches(0.4),
        obs.get("caption", ""),
        font_size=9, italic=True, color=GRAY,
    )

    # Right column — text content
    right_left = Inches(6.5)
    right_width = Inches(6.3)
    y = Inches(0.7)

    # Info grid
    info_items = [
        ("System", obs.get("system", "")),
        ("Component", obs.get("component", "")),
        ("Location", obs.get("location", "")),
        ("Condition", obs.get("condition", "")),
        ("Priority", obs.get("priority", "")),
    ]
    if obs.get("cost_high", 0) > 0:
        cost_str = f"${obs.get('cost_low', 0):,.0f} — ${obs.get('cost_high', 0):,.0f}"
        info_items.append(("Cost Range", cost_str))

    for label, value in info_items:
        _add_text_box(slide, right_left, y, Inches(1.3), Inches(0.25),
                      label + ":", font_size=8, bold=True, color=GRAY)
        _add_text_box(slide, right_left + Inches(1.4), y, Inches(4.8), Inches(0.25),
                      value, font_size=9, color=DARK)
        y += Inches(0.28)

    y += Inches(0.15)

    # Prose
    _add_text_box(slide, right_left, y, right_width, Inches(0.2),
                  "Observation", font_size=10, bold=True, color=NAVY)
    y += Inches(0.3)
    prose_text = obs.get("prose", "")
    # Estimate needed height (rough: ~50 chars per line at 9pt, ~0.18 inches per line)
    prose_lines = max(1, len(prose_text) // 70 + 1)
    prose_height = min(Inches(2.5), Inches(0.2 * prose_lines))
    _add_text_box(slide, right_left, y, right_width, prose_height,
                  prose_text, font_size=9, color=DARK)
    y += prose_height + Inches(0.15)

    # Recommendation
    rec_text = obs.get("recommendation", "")
    if rec_text:
        _add_text_box(slide, right_left, y, right_width, Inches(0.2),
                      "Recommendation", font_size=10, bold=True, color=NAVY)
        y += Inches(0.3)
        rec_lines = max(1, len(rec_text) // 70 + 1)
        rec_height = min(Inches(1.5), Inches(0.2 * rec_lines))
        _add_text_box(slide, right_left, y, right_width, rec_height,
                      rec_text, font_size=9, color=DARK)

    return slide


def _add_funding_summary_slide(prs, observations):
    """Add funding summary table slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
        "Funding Summary", font_size=20, bold=True, color=NAVY,
    )

    # Group costs by system and priority
    system_priority_costs = defaultdict(lambda: defaultdict(float))
    all_priorities = set()

    for obs in observations:
        system = obs.get("system", "Other")
        priority = obs.get("priority", "Unclassified")
        cost_high = obs.get("cost_high", 0)
        all_priorities.add(priority)
        if cost_high > 0:
            system_priority_costs[system][priority] += cost_high

    def priority_sort_key(p):
        pl = p.lower()
        if "immediate" in pl or "1-2" in pl:
            return 0
        if "1-3" in pl or "short" in pl:
            return 1
        if "3-5" in pl or "medium" in pl:
            return 2
        if "6-10" in pl or "long" in pl:
            return 3
        if "10" in pl or "capital" in pl:
            return 4
        return 5

    sorted_priorities = sorted(all_priorities, key=priority_sort_key)
    sorted_systems = sorted(system_priority_costs.keys())

    if not sorted_systems:
        _add_text_box(slide, Inches(0.5), Inches(1.5), Inches(8), Inches(0.5),
                      "No cost data available.", font_size=12, color=GRAY)
        return slide

    # Build table
    cols = len(sorted_priorities) + 2  # System + priorities + Total
    rows = len(sorted_systems) + 2  # header + systems + grand total

    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.0),
                                          Inches(12.3), Inches(0.35) * rows)
    table = table_shape.table

    # Header
    headers = ["System"] + sorted_priorities + ["Total"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data rows
    for i, system in enumerate(sorted_systems):
        row_idx = i + 1
        table.cell(row_idx, 0).text = system
        row_total = 0
        for j, pri in enumerate(sorted_priorities):
            cost = system_priority_costs[system].get(pri, 0)
            cell = table.cell(row_idx, j + 1)
            cell.text = f"${cost:,.0f}" if cost > 0 else "—"
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
            row_total += cost

        total_cell = table.cell(row_idx, cols - 1)
        total_cell.text = f"${row_total:,.0f}" if row_total > 0 else "—"
        for p in total_cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True

        for p in table.cell(row_idx, 0).text_frame.paragraphs:
            p.font.size = Pt(9)

    # Grand total row
    grand_row = len(sorted_systems) + 1
    table.cell(grand_row, 0).text = "TOTAL"
    for p in table.cell(grand_row, 0).text_frame.paragraphs:
        p.font.size = Pt(10)
        p.font.bold = True

    for j, pri in enumerate(sorted_priorities):
        col_total = sum(system_priority_costs[sys].get(pri, 0) for sys in sorted_systems)
        cell = table.cell(grand_row, j + 1)
        cell.text = f"${col_total:,.0f}" if col_total > 0 else "—"
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BG

    grand_total = sum(
        cost for sys_costs in system_priority_costs.values() for cost in sys_costs.values()
    )
    gt_cell = table.cell(grand_row, cols - 1)
    gt_cell.text = f"${grand_total:,.0f}"
    for p in gt_cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
    gt_cell.fill.solid()
    gt_cell.fill.fore_color.rgb = LIGHT_BG

    return slide


def _add_signoff_slide(prs, observations):
    """Add engineer sign-off appendix slide listing all observations with approval status."""
    if not observations:
        return None

    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    _add_text_box(
        slide, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
        "Engineer Sign-Off", font_size=20, bold=True, color=NAVY,
    )

    _add_text_box(
        slide, Inches(0.5), Inches(0.85), Inches(11), Inches(0.4),
        "Observations marked below have been reviewed and approved by the signing engineer.",
        font_size=10, italic=True, color=GRAY,
    )

    cols = 4
    rows = len(observations) + 1
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.4), Inches(12.3), Inches(0.35) * rows
    )
    table = table_shape.table

    col_widths = [Inches(1.2), Inches(4.0), Inches(3.5), Inches(3.6)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ["Obs #", "System", "Approved By", "Date"]
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    for i, obs in enumerate(observations):
        row_idx = i + 1
        approved_by = ""
        approved_at = ""
        if obs.get("approved"):
            approved_by = obs.get("approved_by", "")
            approved_at = obs.get("approved_at", "")
            if approved_at:
                try:
                    approved_at = approved_at[:10]  # ISO date part
                except Exception:
                    pass

        values = [
            obs.get("obs_number", ""),
            obs.get("system", ""),
            approved_by,
            approved_at,
        ]

        for j, val in enumerate(values):
            cell = table.cell(row_idx, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = DARK

            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return slide


def generate_pptx_report(job, job_dir, output_path):
    """Generate the polished PPTX report.

    Args:
        job: Job dict from job.json
        job_dir: Path to the job directory (for images)
        output_path: Path for output .pptx file
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    observations = job.get("processed_observations", [])

    # 1. Cover slide
    _add_cover_slide(prs, job)

    # 2. Observation slides first (we need their indices for hyperlinks)
    # We'll insert command center slides at the beginning after we know observation slide indices
    obs_start_index = 1  # After cover slide, we'll add CC slides, so this will shift

    # Calculate how many CC slides we need
    cc_slide_count = max(1, (len(observations) + ROWS_PER_CC_SLIDE - 1) // ROWS_PER_CC_SLIDE)
    obs_actual_start = 1 + cc_slide_count  # cover + CC slides

    # Build observation slide index map
    obs_slides_map = {}
    for i, obs in enumerate(observations):
        obs_slides_map[obs["obs_number"]] = obs_actual_start + i

    # Add command center slides
    for chunk_idx in range(cc_slide_count):
        start = chunk_idx * ROWS_PER_CC_SLIDE
        end = start + ROWS_PER_CC_SLIDE
        chunk = observations[start:end]
        _add_command_center_slide(prs, chunk, obs_actual_start, obs_slides_map)

    # 3. Observation slides
    for obs in observations:
        _add_observation_slide(prs, obs, job_dir)

    # 4. Funding summary
    _add_funding_summary_slide(prs, observations)

    # 5. Engineer sign-off
    _add_signoff_slide(prs, observations)

    # Save
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
